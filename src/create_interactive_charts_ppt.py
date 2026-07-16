"""Build a PowerPoint review deck from the Final Ver Plotly chart suite."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
CHART_DIR = ROOT / "docs" / "interactive_charts"
OUTPUT = ROOT / "docs" / "PAPL_Final_Ver_Interactive_Chart_Review.pptx"

NAVY = RGBColor(23, 54, 93)
BLUE = RGBColor(47, 85, 151)
ORANGE = RGBColor(245, 166, 35)
LIGHT = RGBColor(242, 246, 251)
MID = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 30, 30)


CHARTS = [
    {
        "file": "01_national_category_opportunity",
        "title": "National Category Opportunity and Parle Share",
        "shows": "Compares total category Sales Value with Parle's calculated value share across the five S4 categories.",
        "insight": "Separates large categories with low Parle capture from categories where Parle already has a stronger position.",
        "method": "Category value: sum of Sales Value (000) where MANUFACTURER is blank. Parle value: sum where MANUFACTURER = PARLE AGRO and ITEM is blank. Share = Parle value ÷ category value.",
    },
    {
        "file": "02_state_opportunity_share",
        "title": "State Category Opportunity and Parle Value Share",
        "shows": "Ranks states by total category Sales Value; bar colour represents Parle's calculated value share.",
        "insight": "Highlights states where a large market pool coincides with low current Parle capture.",
        "method": "Uses the same category-total and Parle-aggregate hierarchy filters as Chart 1, grouped by the state code extracted from Markets.",
    },
    {
        "file": "03_growth_priority_matrix",
        "title": "Growth Prioritization: Opportunity × Current Right-to-Win",
        "shows": "Plots category value against Parle value share. Bubble size represents category value and colour represents the calculated ND gap.",
        "insight": "Supports prioritisation by showing opportunity size, current capture, and distribution headroom together.",
        "method": "Sales calculations follow Charts 1–2. Category ND uses the maximum category-total ND value by state; Parle ND sums Parle aggregate rows. ND gap is clipped at zero.",
    },
    {
        "file": "04_state_category_share_heatmap",
        "title": "Parle Value Share by State and Category",
        "shows": "Displays calculated Parle value share for every available state × S4-category combination.",
        "insight": "Makes category strongholds and weak state-category combinations visible in one view.",
        "method": "For each State and S4CATEGORY, Parle aggregate Sales Value is divided by category-total Sales Value. No PPT values are imported.",
    },
    {
        "file": "05_revenue_states_size_share",
        "title": "Revenue-Led States: Category Size and Share Potential",
        "shows": "Focuses on Gujarat, Telangana, Punjab, and Haryana, comparing category value with calculated Parle share.",
        "insight": "Shows which of the four selected markets combine larger pools with lower current capture.",
        "method": "Uses only Final Ver sales values and the hierarchy filters defined in Chart 1. The four-state selection follows the requested PPT review structure.",
    },
    {
        "file": "06_distribution_headroom",
        "title": "Distribution Gap and Addressable Store Headroom",
        "shows": "Compares average numeric-distribution gap with estimated category-store headroom by state.",
        "insight": "Translates distribution under-penetration into a directional number of addressable stores.",
        "method": "ND gap = category ND − Parle ND by State × S4CATEGORY. Store headroom = ND gap × Number of Stores Retailing ÷ 100. Category ND and stores use maxima; Parle ND sums aggregate rows.",
    },
    {
        "file": "07_pack_format_mix",
        "title": "Parle Sales Mix by Pack Format",
        "shows": "Breaks SKU-level Parle Sales Value into tetra, individual, single-serve, sharing, and family/party formats.",
        "insight": "Shows where current portfolio value is concentrated across consumption and pack-size occasions.",
        "method": "Uses ITEM-populated rows only. Pack type is parsed from Pack Type_Size; size bands follow ≤200 ml, 201–300 ml, 301–600 ml, and >600 ml. Tetra and 80 ml tetra remain separate.",
    },
    {
        "file": "08_pack_mix_by_state",
        "title": "Pack-Format Mix by State",
        "shows": "Compares each state's percentage sales mix across the derived pack-format buckets.",
        "insight": "Reveals states that over-index on family, individual, tetra, or sharing formats.",
        "method": "SKU-level Sales Value is summed by State × Pack Bucket and divided by total SKU-level Sales Value within the state.",
    },
    {
        "file": "09_brand_value",
        "title": "Top Parle Brands by Sales Value",
        "shows": "Ranks the 15 highest-value brands represented in the ITEM-level data.",
        "insight": "Shows portfolio concentration and which brands currently contribute most to recorded value.",
        "method": "Brand is the text before the first backslash in ITEM, after removing any leading #. Sales Value is summed on ITEM-populated rows.",
    },
    {
        "file": "10_portfolio_depth",
        "title": "Portfolio Breadth and Range Depth by State",
        "shows": "Compares the number of active brands and active ITEM descriptions recorded in each state.",
        "insight": "Separates brand breadth from SKU depth and identifies states with limited portfolio activation.",
        "method": "An active record is an ITEM-populated row with Sales Value > 0. Brands and ITEM values are counted distinctly by state.",
    },
    {
        "file": "11_price_ladder",
        "title": "Current Portfolio Value by Price Ladder",
        "shows": "Allocates SKU-level Sales Value across bands defined using Price per Sales Unit.",
        "insight": "Identifies price bands where current recorded portfolio value is relatively concentrated or limited.",
        "method": "Rows require ITEM, Sales Value, and Price per Sales Unit. Bands are ≤₹10, ₹10–15, ₹15–20, ₹20–30, ₹30–50, and ₹50+.",
    },
    {
        "file": "12_distribution_oos",
        "title": "Distribution Reach vs. Out-of-Stock Diagnostic",
        "shows": "Plots state-level Parle numeric distribution against average numeric out-of-stock; bubble size represents Sales Value.",
        "insight": "Helps distinguish markets constrained by reach from those where availability within reached stores may also require attention.",
        "method": "Uses Parle aggregate rows. Numeric Distribution is summed by state, Numeric OOS is averaged, and Sales Value is summed. Treat as a diagnostic because distribution percentages are non-additive.",
    },
]


def add_text(slide, x, y, w, h, text, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title, number=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    prefix = f"{number:02d}  " if number else ""
    add_text(slide, 0.45, 0.12, 12.3, 0.45, prefix + title, 24, WHITE, True)


def add_panel_section(slide, y, heading, body, accent=BLUE):
    add_text(slide, 9.15, y, 3.65, 0.3, heading.upper(), 10, accent, True)
    add_text(slide, 9.15, y + 0.3, 3.65, 1.0, body, 12, BLACK)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    add_text(slide, 0.7, 1.35, 11.9, 1.2, "PAPL Commercial Intelligence", 34, WHITE, True)
    add_text(slide, 0.7, 2.45, 11.9, 0.7, "Interactive Chart Review · Final Ver", 25, RGBColor(180, 210, 240))
    add_text(slide, 0.7, 3.45, 10.7, 1.0, "12 chart-led views built from the populated records in the Final Ver sheet. Click any chart image in this deck to open its interactive Plotly version.", 18, WHITE)
    add_text(slide, 0.7, 6.55, 11.9, 0.35, "Source: PAPL Phase 1- Transposed File_3_vs.xlsx · Apr–Sep 2025", 11, RGBColor(190, 200, 215))

    # Governance slide
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Data and calculation governance")
    items = [
        "Only the Final Ver worksheet is read by analyze_data.py.",
        "The worksheet contains 24,357 populated rows; formatted blank rows are excluded.",
        "No chart values are copied from the previously shared PowerPoint.",
        "Each chart slide states its columns, filters, hierarchy level, and aggregation.",
        "Category totals, Parle aggregate rows, and ITEM-level rows are kept separate.",
        "Derived charts are descriptive; no projected revenue or simulator assumptions are included.",
        "Distribution percentages are non-additive. Charts using summed distribution are explicitly labelled diagnostic and require business validation.",
    ]
    y = 1.15
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background()
        add_text(slide, 1.0, y, 11.4, 0.65, item, 17, BLACK)
        y += 0.78

    # Chart slides
    for number, chart in enumerate(CHARTS, 1):
        slide = prs.slides.add_slide(blank)
        add_header(slide, chart["title"], number)
        image_path = CHART_DIR / f'{chart["file"]}.png'
        html_target = f'interactive_charts/{chart["file"]}.html'
        pic = slide.shapes.add_picture(str(image_path), Inches(0.35), Inches(1.05), width=Inches(8.45), height=Inches(4.83))
        pic.click_action.hyperlink.address = html_target
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.05), Inches(8.45), Inches(4.83))
        border.fill.background(); border.line.color.rgb = RGBColor(210, 218, 228); border.line.width = Pt(1)
        # Transparent border must remain behind the clickable button, so add a dedicated link button.
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(6.1), Inches(3.8), Inches(0.52))
        btn.fill.solid(); btn.fill.fore_color.rgb = BLUE; btn.line.fill.background()
        tf = btn.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "OPEN INTERACTIVE CHART"; r.font.name = "Aptos"; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        btn.click_action.hyperlink.address = html_target
        add_panel_section(slide, 1.05, "What this chart shows", chart["shows"])
        add_panel_section(slide, 2.55, "Business reading", chart["insight"], ORANGE)
        add_panel_section(slide, 4.05, "Final Ver calculation", chart["method"])
        add_text(slide, 0.45, 6.85, 8.1, 0.25, "Click the chart or blue button to open the interactive HTML view.", 10, MID)
        add_text(slide, 9.15, 6.85, 3.7, 0.25, "Source: Final Ver · Apr–Sep 2025", 10, MID, align=PP_ALIGN.RIGHT)

    # Closing slide
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    add_text(slide, 0.8, 1.2, 11.7, 0.8, "Review sequence", 32, WHITE, True)
    add_text(slide, 0.8, 2.1, 11.2, 2.2, "Review each chart with the stated calculation rule. Confirm the intended hierarchy and aggregation—especially for numeric distribution—before treating any result as a final commercial KPI.", 22, WHITE)
    add_text(slide, 0.8, 5.45, 11.2, 0.8, "Interactive files: docs/interactive_charts", 16, RGBColor(180, 210, 240))

    prs.save(OUTPUT)
    print(f"Created {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build_deck()
