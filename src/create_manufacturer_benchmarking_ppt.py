"""Create linked PowerPoint from Original Ver manufacturer charts."""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.util import Inches,Pt

ROOT=Path(__file__).resolve().parents[1]; CHARTS=ROOT/"docs"/"manufacturer_charts"; OUT=ROOT/"docs"/"PAPL_Original_Ver_Manufacturer_Benchmarking_v2.pptx"
NAVY=RGBColor(23,54,93); BLUE=RGBColor(47,85,151); RED=RGBColor(227,30,36); ORANGE=RGBColor(245,166,35); WHITE=RGBColor(255,255,255); BLACK=RGBColor(30,30,30); GREY=RGBColor(100,100,100)
META=[
("01_sales_value_ranking","Manufacturer Sales Value Ranking","Ranks named manufacturers by summed Sales Value for the selected State and S4CATEGORY.","Parle's position, leader, and absolute competitive scale are visible directly from labelled bars.","SUM of Sales Value (000); blank manufacturers excluded."),
("02_value_share","Manufacturer Value Share Comparison","Compares the source Value Share metric across named manufacturers.","Shows whether Parle's sales position is supported by competitive value share.","Unweighted MEAN of Share of Sales Value - Product; no valid source weight exists."),
("03_nd_vs_wd","Numeric Distribution vs Weighted Distribution","Positions manufacturers by numeric and weighted distribution; bubble size is Sales Value.","Separates broad physical reach from concentration in higher-value outlets.","Unweighted MEAN of ND and WD; SUM of Sales Value."),
("04_distribution_productivity","Distribution vs Productivity Matrix","Plots Numeric Distribution against Rate of Sale, sized by Sales Value.","Diagnoses whether Parle trails through reach or productivity.","Unweighted MEAN of ND and ROS; SUM of Sales Value."),
("05_sales_vs_share","Sales Value vs Market Share","Compares absolute sales scale with the source value-share measure.","Highlights manufacturers whose scale and share positions diverge.","SUM of Sales Value; unweighted MEAN of source share."),
("06_store_coverage","Store Coverage Comparison","Ranks manufacturers by Number of Stores Retailing.","Quantifies the physical-store coverage gap between Parle and competitors.","MAX stores at the repeated manufacturer grain to avoid duplicate summation."),
("07_rate_of_sale","Rate of Sale Comparison","Ranks manufacturers by Unweighted ROS.","Shows whether Parle converts its reached stores into sales as productively as competitors.","Unweighted MEAN of ROS."),
("08_price_positioning","Price Positioning","Compares Price per Sales Unit for the selected market/category.","Flags a potential premium or disadvantage; pack-size comparability remains a limitation.","Unweighted MEAN of Price per Sales Unit; interpret only with pack context."),
("09_out_of_stock","Out-of-Stock Comparison","Compares Numeric and Weighted OOS values.","Identifies availability risk where Parle's demand may be constrained by stock-outs.","Unweighted MEAN of Numeric and Weighted OOS."),
("10_performance_heatmap","Manufacturer Performance Heatmap","Normalises six metrics into within-selection percentile scores.","Provides a compact multi-metric benchmark without mixing raw units.","Percentile ranks within the selected State/S4CATEGORY; higher percentile is not favourable for every risk metric."),
("11_opportunity_heatmap","State × S4CATEGORY Opportunity Heatmap","Shows Parle's Sales Value gap to the market leader across combinations.","Darker cells identify larger absolute competitive gaps for review.","Leader Sales Value minus Parle Sales Value; named manufacturers only."),
("12_parle_rank_states","Parle Rank Across States","Tracks Parle Sales Value rank by state for each S4CATEGORY.","Shows where Parle leads and where its competitive rank weakens.","Dense descending rank on summed Sales Value; rank 1 is best."),
("13_leader_gap","Leader Gap Chart","Compares Parle and leader Sales Value for every available state-category combination.","Makes the absolute size of defend-versus-close opportunities explicit.","SUM of Sales Value for named manufacturers; missing Parle combinations excluded from numeric comparison."),
]

def text(slide,x,y,w,h,value,size=14,color=BLACK,bold=False,align=PP_ALIGN.LEFT):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); f=b.text_frame; f.clear(); f.word_wrap=True; f.vertical_anchor=MSO_ANCHOR.TOP; p=f.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=value; r.font.name="Aptos"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; return b
def header(slide,title,n):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.7)); s.fill.solid(); s.fill.fore_color.rgb=NAVY; s.line.fill.background(); text(slide,.4,.12,12.4,.4,f"{n:02d}  {title}",23,WHITE,True)
def panel(slide,y,h,b): text(slide,9.12,y,3.7,.3,h.upper(),10,BLUE,True); text(slide,9.12,y+.28,3.7,1.0,b,12)
def main():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    s=prs.slides.add_slide(blank); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=NAVY; text(s,.7,1.25,12,1,"Original Ver Manufacturer Benchmarking",33,WHITE,True); text(s,.7,2.35,11.5,.8,"Parle Agro versus named manufacturers · 12 states × 5 S4 categories",22,RGBColor(185,210,240)); text(s,.7,3.5,11.5,1.2,"Every chart contains visible value labels. Click a chart or its blue button to open the matching interactive Plotly HTML visual.",18,WHITE); text(s,.7,6.5,11,.3,"Source: Original Ver only · Apr–Sep 2025",11,RGBColor(190,200,215))
    s=prs.slides.add_slide(blank); header(s,"Data governance and limitations",0); bullets=["Original Ver is the only analytical source; Final Ver is not loaded.","Blank-manufacturer rows are separated and never ranked as competitors.","Additive facts use SUM; ratios/rates use documented unweighted means; stores/snapshots use MAX.","The workbook lacks native Facts, CATEGORY, SEGMENT, SUBSEGMENT, BRAND and SUBBRAND columns.","Results describe aggregated retail-audit evidence and do not establish causality."]
    y=1.15
    for b in bullets: text(s,.9,y,11.7,.65,"•  "+b,17); y+=.86
    for n,(slug,title,shows,insight,method) in enumerate(META,1):
        s=prs.slides.add_slide(blank); header(s,title,n); pic=s.shapes.add_picture(str(CHARTS/(slug+".png")),Inches(.3),Inches(1.0),width=Inches(8.5),height=Inches(4.86)); link=f"manufacturer_charts/{slug}.html"; pic.click_action.hyperlink.address=link
        panel(s,1.0,"What the visual shows",shows); panel(s,2.45,"Business interpretation",insight); panel(s,3.9,"Original Ver calculation",method); panel(s,5.35,"Caveat","Use the metric catalogue and selected State/S4CATEGORY context before making an investment decision.")
        btn=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.6),Inches(6.1),Inches(3.9),Inches(.5)); btn.fill.solid(); btn.fill.fore_color.rgb=BLUE; btn.line.fill.background(); tf=btn.text_frame; tf.clear(); tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text="OPEN INTERACTIVE HTML"; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=WHITE; btn.click_action.hyperlink.address=link
        text(s,.4,6.82,8,.25,"Visible labels are embedded; hover provides full precision and manufacturer names.",10,GREY)
    s=prs.slides.add_slide(blank); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=NAVY; text(s,.8,1.3,11.7,.8,"Recommended review workflow",31,WHITE,True); text(s,.8,2.35,11.2,2.2,"1. Select a State and S4CATEGORY in Streamlit.\n2. Confirm Parle rank and leader gap.\n3. Review distribution, productivity, coverage, price, and OOS evidence.\n4. Accept a recommendation only when its triggering metrics are available.",20,WHITE); text(s,.8,6.3,11,.3,"Dashboard: pages/Manufacturer_Comparison.py",13,RGBColor(185,210,240))
    prs.save(OUT); print(f"Created {OUT} with {len(prs.slides)} slides")
if __name__=="__main__": main()
